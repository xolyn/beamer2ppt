import sys
import os

try:
    import fitz
    from pptx import Presentation
    from pptx.util import Inches, Pt
except:
    print("\npackages not installed, make sure you installed `fitz` and `pptx` pakcages!\nuse `pip install python-pptx pymupdf` to install!\n")

def main():    
    if len(sys.argv) != 3:  # Check if the number of arguments is correct
        print("\nInvalid parameter number\nUsage: python beamer2ppt.py [pdf directory] [desired ppt directory]\n")
        return None
    
    pdf_path = str(sys.argv[1])
    pptx_path = str(sys.argv[2])
                
    if not os.path.exists(pdf_path):
        print("\ninvalid pdf file directory!\n")
        return None

    doc = fitz.open(pdf_path)
    pres = Presentation()
    slide_width = pres.slide_width
    slide_height = pres.slide_height
    
    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72))
        img_path = f"page_{page.number}.png"
        pix.save(img_path)
        slide_layout = pres.slide_layouts[6]  # 6: index for a blank slide
        slide = pres.slides.add_slide(slide_layout)
        image_width = Pt(pix.width)
        image_height = Pt(pix.height)
        scale_width = slide_width / image_width
        scale_height = slide_height / image_height
        scale = min(scale_width, scale_height)
        fit_width = image_width * scale
        fit_height = image_height * scale
        
        # centering 
        left = (slide_width - fit_width) / 2
        top = (slide_height - fit_height) / 2

        slide.shapes.add_picture(img_path, left, top, width=fit_width, height=fit_height)
    
    pres.save(pptx_path)
    print(f"Converted {pdf_path} to {pptx_path}")

if __name__ == "__main__":
    main()
